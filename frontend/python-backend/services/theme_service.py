"""
Theme Service - Apply Color and Typography Themes to Documents
Provides comprehensive theming capabilities for DOCX, PPTX, and other documents
"""

from docx import Document
from docx.shared import RGBColor, Pt
from pptx import Presentation
from pptx.util import Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from typing import Dict, List, Any
import io


class ThemeService:
    """Service for applying themes to documents"""

    # Define 10 beautiful theme presets
    THEMES = {
        "ocean": {
            "name": "Ocean Breeze",
            "description": "Cool blues and teals for a calm, professional look",
            "category": "Professional",
            "colors": {
                "primary": "#0077BE",
                "secondary": "#00A3E0",
                "accent": "#00CED1",
                "background": "#F0F8FF",
                "text": "#1C3A52"
            },
            "fonts": {
                "heading": "Calibri",
                "body": "Calibri"
            }
        },
        "sunset": {
            "name": "Sunset Glow",
            "description": "Warm oranges and purples for creative energy",
            "category": "Creative",
            "colors": {
                "primary": "#FF6B35",
                "secondary": "#F7931E",
                "accent": "#C9485B",
                "background": "#FFF5E6",
                "text": "#2D3748"
            },
            "fonts": {
                "heading": "Arial",
                "body": "Arial"
            }
        },
        "forest": {
            "name": "Forest Green",
            "description": "Natural greens for environmental and health topics",
            "category": "Natural",
            "colors": {
                "primary": "#2D5016",
                "secondary": "#4F7942",
                "accent": "#8FBC8F",
                "background": "#F5F9F3",
                "text": "#1A3D23"
            },
            "fonts": {
                "heading": "Georgia",
                "body": "Georgia"
            }
        },
        "corporate": {
            "name": "Corporate Blue",
            "description": "Classic blue and gray for business documents",
            "category": "Professional",
            "colors": {
                "primary": "#003366",
                "secondary": "#0066CC",
                "accent": "#6699CC",
                "background": "#F8F9FA",
                "text": "#212529"
            },
            "fonts": {
                "heading": "Arial",
                "body": "Calibri"
            }
        },
        "vibrant": {
            "name": "Vibrant Pop",
            "description": "Bold colors for attention-grabbing presentations",
            "category": "Creative",
            "colors": {
                "primary": "#FF0080",
                "secondary": "#7928CA",
                "accent": "#FF4785",
                "background": "#FAFAFA",
                "text": "#1A1A1A"
            },
            "fonts": {
                "heading": "Arial",
                "body": "Arial"
            }
        },
        "minimal": {
            "name": "Minimal Gray",
            "description": "Clean grayscale for modern minimalist design",
            "category": "Minimal",
            "colors": {
                "primary": "#2C3E50",
                "secondary": "#34495E",
                "accent": "#7F8C8D",
                "background": "#FFFFFF",
                "text": "#212121"
            },
            "fonts": {
                "heading": "Calibri",
                "body": "Calibri"
            }
        },
        "royal": {
            "name": "Royal Purple",
            "description": "Elegant purples and golds for premium content",
            "category": "Elegant",
            "colors": {
                "primary": "#6B2C91",
                "secondary": "#9B59B6",
                "accent": "#D4AF37",
                "background": "#F9F6FF",
                "text": "#2C1A3F"
            },
            "fonts": {
                "heading": "Georgia",
                "body": "Georgia"
            }
        },
        "tech": {
            "name": "Tech Blue",
            "description": "Modern blue and cyan for technology topics",
            "category": "Technology",
            "colors": {
                "primary": "#0A84FF",
                "secondary": "#00D4FF",
                "accent": "#5E5CE6",
                "background": "#F5F5F7",
                "text": "#1D1D1F"
            },
            "fonts": {
                "heading": "Arial",
                "body": "Arial"
            }
        },
        "warm": {
            "name": "Warm Earth",
            "description": "Earthy browns and oranges for warmth and trust",
            "category": "Natural",
            "colors": {
                "primary": "#8B4513",
                "secondary": "#CD853F",
                "accent": "#DEB887",
                "background": "#FFF8F0",
                "text": "#3E2723"
            },
            "fonts": {
                "heading": "Georgia",
                "body": "Georgia"
            }
        },
        "cool": {
            "name": "Cool Mint",
            "description": "Fresh mint and teal for clean, modern look",
            "category": "Fresh",
            "colors": {
                "primary": "#00A896",
                "secondary": "#02C39A",
                "accent": "#05668D",
                "background": "#F0FFF4",
                "text": "#1A4D2E"
            },
            "fonts": {
                "heading": "Calibri",
                "body": "Calibri"
            }
        }
    }

    @staticmethod
    def get_all_themes() -> List[Dict[str, Any]]:
        """Get all available themes with metadata"""
        return [
            {
                "id": theme_id,
                **theme_data
            }
            for theme_id, theme_data in ThemeService.THEMES.items()
        ]

    @staticmethod
    def get_theme(theme_id: str) -> Dict[str, Any]:
        """Get a specific theme by ID"""
        if theme_id not in ThemeService.THEMES:
            raise ValueError(f"Theme '{theme_id}' not found")
        return {
            "id": theme_id,
            **ThemeService.THEMES[theme_id]
        }

    @staticmethod
    def apply_theme_to_docx(doc: Document, theme_id: str) -> Document:
        """
        Apply a theme to a DOCX document

        Args:
            doc: Document object
            theme_id: Theme identifier

        Returns:
            Modified document
        """
        theme = ThemeService.THEMES.get(theme_id)
        if not theme:
            raise ValueError(f"Theme '{theme_id}' not found")

        # Extract colors (convert hex to RGB)
        primary = ThemeService._hex_to_rgb(theme["colors"]["primary"])
        secondary = ThemeService._hex_to_rgb(theme["colors"]["secondary"])
        text_color = ThemeService._hex_to_rgb(theme["colors"]["text"])

        # Apply theme to styles
        styles = doc.styles

        # Update heading styles
        if "Heading 1" in styles:
            h1 = styles["Heading 1"]
            h1.font.name = theme["fonts"]["heading"]
            h1.font.color.rgb = RGBColor(*primary)
            h1.font.size = Pt(24)

        if "Heading 2" in styles:
            h2 = styles["Heading 2"]
            h2.font.name = theme["fonts"]["heading"]
            h2.font.color.rgb = RGBColor(*secondary)
            h2.font.size = Pt(18)

        if "Heading 3" in styles:
            h3 = styles["Heading 3"]
            h3.font.name = theme["fonts"]["heading"]
            h3.font.color.rgb = RGBColor(*secondary)
            h3.font.size = Pt(14)

        # Update normal text style
        if "Normal" in styles:
            normal = styles["Normal"]
            normal.font.name = theme["fonts"]["body"]
            normal.font.color.rgb = RGBColor(*text_color)
            normal.font.size = Pt(11)

        return doc

    @staticmethod
    def apply_theme_to_pptx(prs: Presentation, theme_id: str) -> Presentation:
        """
        Apply a theme to a PPTX presentation

        Args:
            prs: Presentation object
            theme_id: Theme identifier

        Returns:
            Modified presentation
        """
        theme = ThemeService.THEMES.get(theme_id)
        if not theme:
            raise ValueError(f"Theme '{theme_id}' not found")

        # Extract colors
        primary = ThemeService._hex_to_rgb(theme["colors"]["primary"])
        text_color = ThemeService._hex_to_rgb(theme["colors"]["text"])

        # Apply theme to all slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = theme["fonts"]["body"]
                            if shape == slide.shapes.title:
                                run.font.color.rgb = PptxRGBColor(*primary)
                                run.font.size = PptxPt(32)
                            else:
                                run.font.color.rgb = PptxRGBColor(*text_color)

        return prs

    @staticmethod
    def _hex_to_rgb(hex_color: str) -> tuple:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    @staticmethod
    def save_docx_to_bytes(doc: Document) -> bytes:
        """Save DOCX document to bytes"""
        bytes_io = io.BytesIO()
        doc.save(bytes_io)
        bytes_io.seek(0)
        return bytes_io.read()

    @staticmethod
    def save_pptx_to_bytes(prs: Presentation) -> bytes:
        """Save PPTX presentation to bytes"""
        bytes_io = io.BytesIO()
        prs.save(bytes_io)
        bytes_io.seek(0)
        return bytes_io.read()
