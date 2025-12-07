"""
PPTX Service - PowerPoint Presentation Generation and Manipulation
Provides comprehensive presentation creation and editing capabilities
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Dict, Optional
import io


class PPTXService:
    """Service for creating and manipulating PowerPoint presentations"""

    @staticmethod
    def create_blank_presentation() -> Presentation:
        """Create a new blank presentation"""
        return Presentation()

    @staticmethod
    def create_presentation_from_template(template_type: str) -> Presentation:
        """
        Create a presentation from a predefined template

        Args:
            template_type: Type of template (business, pitch, report, etc.)

        Returns:
            Presentation object
        """
        prs = Presentation()

        if template_type == "business":
            PPTXService._add_business_template(prs)
        elif template_type == "pitch":
            PPTXService._add_pitch_template(prs)
        elif template_type == "report":
            PPTXService._add_report_template(prs)
        elif template_type == "educational":
            PPTXService._add_educational_template(prs)
        else:
            # Default blank template
            pass

        return prs

    @staticmethod
    def _add_business_template(prs: Presentation):
        """Add business presentation template slides"""
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Business Presentation"
        subtitle.text = "Your Company Name"

        # Content slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "Key Points"
        tf = body.text_frame
        tf.text = "Strategic Overview"

        p = tf.add_paragraph()
        p.text = "Financial Performance"
        p.level = 0

        p = tf.add_paragraph()
        p.text = "Future Outlook"
        p.level = 0

    @staticmethod
    def _add_pitch_template(prs: Presentation):
        """Add investor pitch template slides"""
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Investment Pitch"
        subtitle.text = "Company Vision"

        # Problem slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "The Problem"
        tf = body.text_frame
        tf.text = "Market pain point we're solving"

    @staticmethod
    def _add_report_template(prs: Presentation):
        """Add report template slides"""
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Quarterly Report"
        subtitle.text = "Q4 2024"

        # Executive Summary
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        body = slide.placeholders[1]

        title.text = "Executive Summary"
        tf = body.text_frame
        tf.text = "Key findings and recommendations"

    @staticmethod
    def _add_educational_template(prs: Presentation):
        """Add educational template slides"""
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Course Title"
        subtitle.text = "Lesson Overview"

    @staticmethod
    def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> Presentation:
        """Add a title slide to the presentation"""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        return prs

    @staticmethod
    def add_content_slide(prs: Presentation, title: str, content: List[str]) -> Presentation:
        """Add a content slide with bullet points"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = content[0] if content else ""

        for item in content[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

        return prs

    @staticmethod
    def add_image_slide(prs: Presentation, title: str, image_path: str) -> Presentation:
        """Add a slide with an image"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title

        # Add image
        left = Inches(1)
        top = Inches(2)
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(8))

        return prs

    @staticmethod
    def add_table_slide(prs: Presentation, title: str, data: List[List[str]]) -> Presentation:
        """Add a slide with a table"""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        slide.shapes.title.text = title

        # Add table
        rows = len(data)
        cols = len(data[0]) if data else 0

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Populate table
        for i, row_data in enumerate(data):
            for j, cell_data in enumerate(row_data):
                table.cell(i, j).text = str(cell_data)

        return prs

    @staticmethod
    def save_presentation(prs: Presentation, output_path: str):
        """Save presentation to file"""
        prs.save(output_path)

    @staticmethod
    def save_presentation_to_bytes(prs: Presentation) -> bytes:
        """Save presentation to bytes for download"""
        bytes_io = io.BytesIO()
        prs.save(bytes_io)
        bytes_io.seek(0)
        return bytes_io.read()

    @staticmethod
    def get_template_types() -> List[Dict[str, str]]:
        """Get available presentation templates"""
        return [
            {
                "id": "blank",
                "name": "Blank Presentation",
                "description": "Start with a blank canvas"
            },
            {
                "id": "business",
                "name": "Business Presentation",
                "description": "Professional business presentation template"
            },
            {
                "id": "pitch",
                "name": "Investor Pitch",
                "description": "Template for startup pitch decks"
            },
            {
                "id": "report",
                "name": "Report Presentation",
                "description": "Quarterly or annual report template"
            },
            {
                "id": "educational",
                "name": "Educational",
                "description": "Teaching and training presentation template"
            }
        ]

    @staticmethod
    def apply_theme(prs: Presentation, theme_name: str) -> Presentation:
        """
        Apply a color theme to the presentation

        Args:
            prs: Presentation object
            theme_name: Theme name (corporate, vibrant, minimal, etc.)

        Returns:
            Modified presentation
        """
        themes = {
            "corporate": {
                "title_color": RGBColor(0, 66, 124),
                "text_color": RGBColor(51, 51, 51),
                "accent_color": RGBColor(255, 102, 0)
            },
            "vibrant": {
                "title_color": RGBColor(255, 51, 102),
                "text_color": RGBColor(34, 34, 34),
                "accent_color": RGBColor(102, 204, 255)
            },
            "minimal": {
                "title_color": RGBColor(34, 34, 34),
                "text_color": RGBColor(85, 85, 85),
                "accent_color": RGBColor(153, 153, 153)
            }
        }

        if theme_name not in themes:
            return prs

        theme = themes[theme_name]

        # Apply theme to all slides
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if shape == slide.shapes.title:
                                run.font.color.rgb = theme["title_color"]
                            else:
                                run.font.color.rgb = theme["text_color"]

        return prs
